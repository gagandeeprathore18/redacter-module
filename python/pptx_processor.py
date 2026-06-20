import pptx
import re
from PIL import Image
import io
import os
import json
from redact_engine import (
    redact_text, is_logo_match, redact_paragraph_runs,
    redact_paragraph_runs_with_pattern,
    TABLE_ROW_NAME_KEYWORD_PATTERN,
    SAFE_COLUMN_HEADER_PATTERN,
    TABLE_ROW_LOCATION_KEYWORD_PATTERN, SUBMISSION_LOCATION_PATTERN
)
from redaction.stop_patterns import should_stop_block

from redaction.table_parser import normalize_table
from redaction.target_detector import detect_targets
from redaction.relationship_detector import get_block_coordinates
from redaction.block_extractor import extract_block_cells
from redaction.block_redactor import redact_cell
from redaction.validator import validate_table_redaction
from redaction.protected_zone_detector import reset_protected_zone, set_protected_zone

def process_table(table):
    grid = normalize_table(table)
    target_coords = detect_targets(grid)
    all_block_coords = set()
    for coord in target_coords:
        block = get_block_coordinates(grid, coord)
        all_block_coords.update(block)
        
    if all_block_coords:
        cells_to_redact = extract_block_cells(grid, all_block_coords)
        for cell in cells_to_redact:
            redact_cell(cell, is_docx=False)
        final_grid = normalize_table(table)
        validate_table_redaction(grid, final_grid, all_block_coords)

    # Get column headers
    headers = []
    if len(table.rows) > 0:
        for cell in table.rows[0].cells:
            headers.append(cell.text_frame.text.strip() if cell.text_frame else "")

    for row_idx, row in enumerate(table.rows):
        # Keyword scanning
        row_has_name_keyword = False
        row_has_location_keyword = False
        label_cell_idx = -1

        # Label cell text (first column) for continuation detection
        label_cell_text = ""
        if row.cells and row.cells[0].text_frame:
            label_cell_text = row.cells[0].text_frame.text.strip()

        for idx, cell in enumerate(row.cells):
            if (row_idx, idx) in all_block_coords:
                continue
            if cell.text_frame:
                cell_text = cell.text_frame.text.strip()
                if TABLE_ROW_NAME_KEYWORD_PATTERN.search(cell_text):
                    row_has_name_keyword = True
                    label_cell_idx = idx
                if TABLE_ROW_LOCATION_KEYWORD_PATTERN.search(cell_text):
                    row_has_location_keyword = True

        for idx, cell in enumerate(row.cells):
            if (row_idx, idx) in all_block_coords:
                continue
            if cell.text_frame:
                cell_text = cell.text_frame.text.strip()
                set_protected_zone(cell_text)

                # Submission location rows: blank ALL cells entirely
                if row_has_location_keyword:
                    for p in cell.text_frame.paragraphs:
                        set_protected_zone(p.text)
                        for run in p.runs:
                            run.text = " "
                    continue

                redact_names_in_cell = row_has_name_keyword and (idx != label_cell_idx)
                for p in cell.text_frame.paragraphs:
                    set_protected_zone(p.text)
                    redact_paragraph_runs(p.runs, redact_all_names=redact_names_in_cell)

def should_stop_pptx_block(text: str) -> bool:
    """Delegates to the shared stop-block checker."""
    return should_stop_block(text)

def get_blank_image_bytes() -> bytes:
    img = Image.new("RGBA", (1, 1), (255, 255, 255, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

def process_shapes(shapes):
    for shape in shapes:
        # 1. Text Frame
        if shape.has_text_frame:
            in_location_block = False
            for p in shape.text_frame.paragraphs:
                set_protected_zone(p.text)
                para_text = p.text.strip()
                if not para_text:
                    if in_location_block:
                        for run in p.runs:
                            run.text = " "
                    continue
                if SUBMISSION_LOCATION_PATTERN.search(para_text) or TABLE_ROW_LOCATION_KEYWORD_PATTERN.search(para_text):
                    for run in p.runs:
                        run.text = " "
                    in_location_block = True
                elif in_location_block:
                    if should_stop_pptx_block(para_text):
                        in_location_block = False
                        redact_paragraph_runs(p.runs)
                    else:
                        for run in p.runs:
                            run.text = " "
                else:
                    redact_paragraph_runs(p.runs)
                        
        # 2. Table
        if shape.has_table:
            process_table(shape.table)
                                    
        # 3. Group Shapes (recursive)
        if hasattr(shape, "shapes"):
            process_shapes(shape.shapes)

def process_pptx(input_path: str, output_path: str):
    os.environ["CURRENT_DOCUMENT_ID"] = os.path.basename(input_path)
    from redaction.ownership_manager import clear_issuing_university, determine_issuing_university
    from redaction.redaction_debug_logger import set_document_context
    clear_issuing_university()
    set_document_context(document=os.path.basename(input_path), source="pptx")
    
    prs = pptx.Presentation(input_path)
    
    # 1. Gather image metadata from slides/layouts/masters to run OCR detection first
    image_metadata_map = {} # hash -> list of dicts (page/slide_index, location, width, height, bytes)
    image_counts = {}       # hash -> repeat_count
    
    def check_shapes(shapes, is_master_or_layout, slide_index):
        for shape in shapes:
            # Check for PICTURE shape type (type 13)
            if shape.shape_type == 13: # MSO_SHAPE_TYPE.PICTURE
                try:
                    img = shape.image
                    img_hash = hash(img.blob)
                    width_px = int(shape.width / 9525) if shape.width else 0
                    height_px = int(shape.height / 9525) if shape.height else 0
                    
                    location = "body"
                    if is_master_or_layout:
                        location = "header"
                    else:
                        # Slide width and height are available on prs.slide_width / slide_height
                        # e.g., if top is in top 15%
                        if shape.top < prs.slide_height * 0.15:
                            location = "header"
                        elif shape.top > prs.slide_height * 0.85:
                            location = "footer"
                            
                    meta = {
                        "page": slide_index + 1,
                        "location": location,
                        "width": width_px,
                        "height": height_px,
                        "bytes": img.blob
                    }
                    if img_hash not in image_metadata_map:
                        image_metadata_map[img_hash] = []
                    image_metadata_map[img_hash].append(meta)
                    image_counts[img_hash] = image_counts.get(img_hash, 0) + 1
                except Exception:
                    pass
            if hasattr(shape, "shapes"):
                check_shapes(shape.shapes, is_master_or_layout, slide_index)
                
    # Scan slide layouts and masters
    for layout in prs.slide_layouts:
        check_shapes(layout.shapes, True, 0)
    for master in prs.slide_masters:
        check_shapes(master.shapes, True, 0)
        
    # Scan all slides
    for idx, slide in enumerate(prs.slides):
        check_shapes(slide.shapes, False, idx)
 
    # 2. Evaluate branding decision & Determine Issuing University
    images_to_redact = set()
    try:
        import sys
        root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        if root_dir not in sys.path:
            sys.path.insert(0, root_dir)
            
        from branding.image_preprocessor import preprocess_image
        from branding.ocr_detector import OCRDetector
        from branding.branding_decision import BrandingDecisionEngine
        
        ocr = OCRDetector()
        decision_engine = BrandingDecisionEngine()
        
        for img_hash, occurrences in image_metadata_map.items():
            first_occ = occurrences[0]
            img_meta = {
                "location": first_occ["location"],
                "page": first_occ["page"],
                "width": first_occ["width"],
                "height": first_occ["height"],
                "repeat_count": image_counts[img_hash]
            }
            
            preprocessed_bytes = preprocess_image(first_occ["bytes"])
            ocr_text = ocr.extract_text(preprocessed_bytes)
            
            should_remove, score, breakdown = decision_engine.evaluate_image(img_meta, ocr_text)
            print(f"PPTX Evaluated image: score={score}, ocr='{ocr_text}', decision={'REMOVE' if should_remove else 'KEEP'}, breakdown={breakdown}")
            
            # Determine issuing university from logo OCR text
            determine_issuing_university(ocr_text)
            
            if should_remove or is_logo_match(first_occ["bytes"]):
                images_to_redact.add(img_hash)
    except Exception as e:
        print(f"Error during PPTX branding detection: {e}")
        # Fallback to simple is_logo_match
        for img_hash, occurrences in image_metadata_map.items():
            if is_logo_match(occurrences[0]["bytes"]):
                images_to_redact.add(img_hash)

    # 2.5. Scan pass for LLM Escalations
    from redaction.escalation_manager import clear_cache, register_candidate_scan, run_gpt_review, get_document_metrics
    from redaction.entity_classifier import classify_entity
    from redaction.ownership_manager import get_issuing_university
    
    clear_cache()
    reset_protected_zone()
    
    def scan_element_text(text, context="", in_table=False, is_template=False):
        if not text or not text.strip():
            return
        if is_template:
            from redaction.ownership_manager import scan_text_for_universities
            scan_text_for_universities(text)
        set_protected_zone(text)
        classification, action, reasons, score = classify_entity(text, context=context, in_table=in_table)
        register_candidate_scan(text, context, classification, action, score, reasons)

        # Scan for Proximity Names
        from redact_engine import NAME_PROXIMITY_PATTERN, NAME_PATTERN, TABLE_ROW_NAME_KEYWORD_PATTERN
        for match in NAME_PROXIMITY_PATTERN.finditer(text):
            name_str = match.group(1).strip()
            if name_str and len(name_str) > 2:
                classification, action, reasons, score = classify_entity(name_str, context=match.group(0), in_table=in_table)
                register_candidate_scan(name_str, match.group(0), classification, action, score, reasons)
                
        # Scan for general Names if name keyword exists
        if TABLE_ROW_NAME_KEYWORD_PATTERN.search(text):
            for match in NAME_PATTERN.finditer(text):
                name_str = match.group(0).strip()
                if name_str and len(name_str) > 2:
                    start_idx = max(0, match.start() - 120)
                    end_idx = min(len(text), match.end() + 120)
                    ctx = text[start_idx:end_idx]
                    classification, action, reasons, score = classify_entity(name_str, context=ctx, in_table=in_table)
                    register_candidate_scan(name_str, ctx, classification, action, score, reasons)

        # Scan for Date Candidates
        from redaction.date_time_detector import find_date_time_spans
        for start, end, matched_str, m_type in find_date_time_spans(text):
            if m_type == "DATE":
                matched_str = matched_str.strip()
                if matched_str and len(matched_str) > 2:
                    start_idx = max(0, start - 120)
                    end_idx = min(len(text), end + 120)
                    ctx = text[start_idx:end_idx]
                    classification, action, reasons, score = classify_entity(
                        matched_str, 
                        context=ctx, 
                        source_detector="DATE_CANDIDATE_PATTERN",
                        in_table=in_table
                    )
                    register_candidate_scan(matched_str, ctx, classification, action, score, reasons)

    def scan_shapes(shapes, is_template=False):
        for shape in shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    set_protected_zone(p.text)
                    scan_element_text(p.text, is_template=is_template)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame:
                            set_protected_zone(cell.text_frame.text)
                            scan_element_text(cell.text_frame.text, in_table=True, is_template=is_template)
                            for p in cell.text_frame.paragraphs:
                                set_protected_zone(p.text)
                                scan_element_text(p.text, in_table=True, is_template=is_template)
            if hasattr(shape, "shapes"):
                scan_shapes(shape.shapes, is_template=is_template)

    # Scan slides
    for slide in prs.slides:
        scan_shapes(slide.shapes, is_template=False)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for p in slide.notes_slide.notes_text_frame.paragraphs:
                scan_element_text(p.text, is_template=False)
                
    # Scan layouts and masters
    for layout in prs.slide_layouts:
        scan_shapes(layout.shapes, is_template=True)
    for master in prs.slide_masters:
        scan_shapes(master.shapes, is_template=True)

    # Run the GPT batch review if there are escalated candidates
    issuing_univ = get_issuing_university()
    doc_id = os.path.basename(input_path)
    run_gpt_review(issuing_university=issuing_univ, doc_id=doc_id)
    
    # Log passive telemetry metrics
    metrics = get_document_metrics(doc_id)
    print(f"Hybrid classification metrics: {json.dumps(metrics)}")

    # 3. Redact text in all slides (now that ISSUING_UNIVERSITY is determined)
    for slide in prs.slides:
        process_shapes(slide.shapes)
        # Process slide notes if present
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            for p in slide.notes_slide.notes_text_frame.paragraphs:
                redact_paragraph_runs(p.runs)
                        
    # 4. Redact text in slide layouts and master slides to catch template-level texts
    for layout in prs.slide_layouts:
        process_shapes(layout.shapes)
    for master in prs.slide_masters:
        process_shapes(master.shapes)
 
    # 5. Redact Logo Images at the package level
    blank_img_bytes = get_blank_image_bytes()
    for part in prs.part.package.iter_parts():
        if part.content_type and "image" in part.content_type:
            try:
                blob = part._blob
                blob_hash = hash(blob)
                if blob_hash in images_to_redact or is_logo_match(blob):
                    part._blob = blank_img_bytes

                    # Phase 7 Log: Final Decision
                    try:
                        from redaction.redaction_audit import RedactionAudit
                        RedactionAudit.log({
                            "candidate": f"[IMAGE hash={blob_hash}]",
                            "stage": "FINAL_DECISION",
                            "classification": "UNIVERSITY_BRANDING",
                            "decision": "REDACT",
                            "decision_source": "UNIVERSITY_BRANDING"
                        })
                    except Exception:
                        pass
                    # Phase 8 Log: Redaction Geometry
                    try:
                        from redaction.redaction_audit import RedactionAudit
                        RedactionAudit.log({
                            "candidate": f"[IMAGE hash={blob_hash}]",
                            "page": None,
                            "stage": "REDACTION_GEOMETRY",
                            "bbox": None,
                            "bbox_width": None,
                            "bbox_height": None,
                            "ocr_text_inside_bbox": f"[IMAGE hash={blob_hash}]"
                        })
                    except Exception:
                        pass
            except Exception:
                pass
    prs.save(output_path)
    
    # Write final audit summary report at the very end
    try:
        from redaction.redaction_audit import RedactionAudit
        from redaction.escalation_manager import LOGS_DIR
        summary = RedactionAudit.generate_summary(doc_id)
        summary_file = os.path.join(LOGS_DIR, "audit_summary.json")
        with open(summary_file, "w", encoding="utf-8") as sf:
            json.dump(summary, sf, indent=2)
    except Exception as re_err:
        print(f"Error generating final redaction audit report: {re_err}")

