import os
import sys
import docx
import fitz
import pptx
from redact_engine import redact_text

# Add root directory to python path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from python.docx_processor import process_docx
from python.pdf_processor import process_pdf
from python.pptx_processor import process_pptx

# Disable real GPT calling during test_redactor run by clearing OPENAI_API_KEY
import os
os.environ["TESTING"] = "true"
if "OPENAI_API_KEY" in os.environ:
    del os.environ["OPENAI_API_KEY"]

def create_test_docx(filename: str):
    doc = docx.Document()
    doc.add_heading("University Student Profile", level=1)
    
    doc.add_paragraph("Student Identification Number: ST12345")
    doc.add_paragraph("Registration Code: REG2024001")
    
    # Styled split runs to mimic user document
    p_split = doc.add_paragraph("Submission time and date: ")
    r1 = p_split.add_run("14:00, ")
    r1.bold = True
    r1.font.color.rgb = docx.shared.RGBColor(255, 0, 0) # red color
    p_split.add_run("19")
    r3 = p_split.add_run("th")
    r3.font.superscript = True
    p_split.add_run(" June 2026")
    
    # Table test
    table = doc.add_table(rows=15, cols=2)
    table.cell(0, 0).paragraphs[0].text = "Field"
    table.cell(0, 1).paragraphs[0].text = "Detail"
    
    table.cell(1, 0).paragraphs[0].text = "Contact Email"
    table.cell(1, 1).paragraphs[0].text = "admissions@harvard.edu"
    
    table.cell(2, 0).paragraphs[0].text = "Support Phone"
    table.cell(2, 1).paragraphs[0].text = "+1 (123) 456-7890"
    
    # Add plain non-issuing university domain to verify preservation
    doc.add_paragraph("Website: www.harvard.edu/admissions")

    # Add adjacent cell key-value test cases
    table.cell(3, 0).paragraphs[0].text = "Submission time and date:"
    p_val1 = table.cell(3, 1).paragraphs[0]
    r_val1_1 = p_val1.add_run("14:00, ")
    r_val1_1.bold = True
    p_val1.add_run("19")
    r_val1_3 = p_val1.add_run("th")
    r_val1_3.font.superscript = True
    p_val1.add_run(" June 2026")

    table.cell(4, 0).paragraphs[0].text = "Target feedback time and date:"
    p_val2 = table.cell(4, 1).paragraphs[0]
    r_val2_1 = p_val2.add_run("14:00, ")
    r_val2_1.bold = True
    p_val2.add_run("10")
    r_val2_3 = p_val2.add_run("th")
    r_val2_3.font.superscript = True
    p_val2.add_run(" July 2026")

    # Add resit and military time table row test cases (mimicking user screenshot)
    table.cell(5, 0).paragraphs[0].text = "Resit date"
    table.cell(5, 1).paragraphs[0].text = "21st August 2026 (For both Group presentation and or Individual Reflective statement)"

    table.cell(6, 0).paragraphs[0].text = "Submission due dates/time"
    table.cell(6, 1).paragraphs[0].text = "Individual group members should submit their Option A Part 2 document via Turnitin on the VLE NOT LATER than 4:00pm/1600hrs on 12 June 2026"

    # Table row test for names
    table.cell(7, 0).paragraphs[0].text = "Tutor Name"
    table.cell(7, 1).paragraphs[0].text = "John Smith"

    # Add Programme and Module name rows to verify they are NOT redacted
    table.cell(8, 0).paragraphs[0].text = "Programme:"
    table.cell(8, 1).paragraphs[0].text = "BSc Health and Social Care"

    table.cell(9, 0).paragraphs[0].text = "Module name:"
    table.cell(9, 1).paragraphs[0].text = "Research Methods in Social Sciences"

    # Add Draft, Formative, and parenthetical submission rows
    table.cell(10, 0).paragraphs[0].text = "Draft Submission"
    table.cell(10, 1).paragraphs[0].text = "25th October 2026"

    table.cell(11, 0).paragraphs[0].text = "Draft Submission (Mandatory)"
    table.cell(11, 1).paragraphs[0].text = "26th October 2026"

    table.cell(12, 0).paragraphs[0].text = "Formative Submission"
    table.cell(12, 1).paragraphs[0].text = "27th October 2026"

    table.cell(13, 0).paragraphs[0].text = "Date and Time of Submission"
    table.cell(13, 1).paragraphs[0].text = "29th October 2026"

    table.cell(14, 0).paragraphs[0].text = "Submission Date & Time"
    table.cell(14, 1).paragraphs[0].text = "30th October 2026"

    # Paragraph tests for names and parenthetical submission dates
    doc.add_paragraph("Draft Submission (Mandatory): 28th October 2026")
    doc.add_paragraph("Date and Time of Submission: 31st October 2026")
    doc.add_paragraph("Module leader: Adeeba Ahmad")
    doc.add_paragraph("Name/Signed: Nargisa Simansone")
    doc.add_paragraph("Module lead's name: John Connor")
    doc.add_paragraph("Module Lead: Sarah Connor")
    doc.add_paragraph("Academic Year 2025-26")
    doc.add_paragraph("Module Code: BM414")


    # Submission location tests
    doc.add_paragraph("Submission location: Turnitin VLE online portal")
    doc.add_paragraph("Submit to: Blackboard submission folder")
    
    # Free-text multi-paragraph options block
    doc.add_paragraph("Submission Location:")
    doc.add_paragraph("Option A: Online Upload")
    doc.add_paragraph("Option B: In-person Hand-in")
    doc.add_paragraph("Tutor Name: Prof. Green")
    
    # Fuzzy matching / OCR error test case
    doc.add_paragraph("Submisslon Locatlon:")
    doc.add_paragraph("Option A: OCR Failure Test")

    # Exclusions preservation tests (headings and book titles)
    doc.add_paragraph("Recommended Reading:")
    doc.add_paragraph("Managing Innovation by Tidd and Bessant")
    doc.add_paragraph("Reference List:")
    doc.add_paragraph("1. Tidd, J. and Bessant, J. (2021) Managing Innovation: Integrating Technological, Market and Organisational Change. 7th edn. Hoboken: Wiley.")
    doc.add_paragraph("Learning Outcomes: LO1, LO2, LO3, LO4")
    doc.add_paragraph("Module Guide: BM304 Module Guide")
    doc.add_paragraph("Assessment Brief: NCG Assessment Brief")
    doc.add_paragraph("Confidentiality:")
    doc.add_paragraph("This module handbook contains confidential information regarding the curriculum.")
    doc.add_paragraph("Academic Integrity:")
    doc.add_paragraph("Students must maintain academic integrity at all times.")

    # Add required human name candidates
    doc.add_paragraph("Tutor Name: Claire Ngo")
    doc.add_paragraph("Assessor: Sarah Johnson")
    doc.add_paragraph("Internal Verifier: Michael Brown")

    # Add BNU / Buckinghamshire references to verify they get redacted (since logo is university_1/BNU)
    doc.add_paragraph("This is a BNU course module hand-out.")
    doc.add_paragraph("Please contact helpdesk@bucks.ac.uk for support.")
    doc.add_paragraph("Issued by Buckinghamshire New University.")

    # Header and Footer
    section = doc.sections[0]
    header = section.header
    header.paragraphs[0].text = "Official Harvard Document - CONFIDENTIAL"
    footer = section.footer
    footer.paragraphs[0].text = "Drafted by Harvard Admits. Postcode: SW1A 1AA"

    # Layout Type 1 Table
    doc.add_heading("Layout Type 1 Table", level=2)
    t1 = doc.add_table(rows=1, cols=2)
    t1.cell(0, 0).paragraphs[0].text = "Submission Location"
    t1.cell(0, 1).paragraphs[0].text = "Turnitin Submission Portal"

    # Layout Type 2 Table
    doc.add_heading("Layout Type 2 Table", level=2)
    t2 = doc.add_table(rows=3, cols=2)
    t2.cell(0, 0).paragraphs[0].text = "Submission Location"
    t2.cell(0, 1).paragraphs[0].text = "Option A: Presentation"
    t2.cell(1, 0).paragraphs[0].text = "" # Continuation Row
    t2.cell(1, 1).paragraphs[0].text = "Option B: Turnitin Submission"
    t2.cell(2, 0).paragraphs[0].text = "Feedback Date" # Stop Condition Row
    t2.cell(2, 1).paragraphs[0].text = "10 July 2026"

    # Layout Type 3 Table (Vertical)
    doc.add_heading("Layout Type 3 Table", level=2)
    t3 = doc.add_table(rows=2, cols=1)
    t3.cell(0, 0).paragraphs[0].text = "Submission Location"
    t3.cell(1, 0).paragraphs[0].text = "Turnitin Submission Portal"

    # Layout Type 4 Table (Merged Cell)
    doc.add_heading("Layout Type 4 Table", level=2)
    t4 = doc.add_table(rows=2, cols=2)
    t4.cell(0, 0).paragraphs[0].text = "Submission Location"
    # Merge row 0 col 1 and row 1 col 1
    merged_cell = t4.cell(0, 1).merge(t4.cell(1, 1))
    merged_cell.paragraphs[0].text = "Option A + Option B + Option C"

    # Insert a logo that should be redacted
    logo_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "logos", "university_1.png")
    if os.path.exists(logo_path):
        doc.add_picture(logo_path)
        
    doc.save(filename)
    print(f"Created test DOCX: {filename}")

def create_test_pdf(filename: str):
    doc = fitz.open()
    page = doc.new_page()
    
    # Header Zone elements (Top 20%, y < 168.4)
    # 1. Header logo
    logo_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "logos", "university_2.png")
    if os.path.exists(logo_path):
        page.insert_image(fitz.Rect(450, 10, 550, 45), filename=logo_path)
        
    # 2. Header address block (multiline, road keyword, postal code, campus keyword)
    address_text = "Stanford University Campus\n100 Stanford Road, CA 94305"
    page.insert_textbox(fitz.Rect(200, 10, 400, 60), address_text, fontsize=10)
    
    # 3. Header contact block
    page.insert_text((50, 50), "Office of Admissions - Stanford University", fontsize=16)
    page.insert_text((50, 100), "Applicant Reference ID: APP99881", fontsize=12)
    page.insert_text((50, 130), "Contact us at registry@stanford.edu or call 650-123-4567", fontsize=12)
    page.insert_text((50, 160), "Please send replies to Stanford, CA 94305, USA", fontsize=12)
    page.insert_text((50, 190), "Website: www.stanford.edu/admissions", fontsize=12)
    page.insert_text((50, 210), "Target feedback time: 14:30 on 12 December 2025", fontsize=12)
    page.insert_text((50, 230), "Resit date: 21st August 2026 (For both Group presentation and or Individual Reflective statement)", fontsize=10)
    page.insert_text((50, 250), "Submission due dates/time: Submit document NOT LATER than 4:00pm/1600hrs on 12 June 2026", fontsize=10)
    
    page.insert_text((50, 270), "Name/Signed: Nargisa Simansone", fontsize=10)
    page.insert_text((50, 290), "Module leader: Adeeba Ahmad", fontsize=10)
    page.insert_text((50, 310), "Programme: BSc Health and Social Care", fontsize=10)
    page.insert_text((50, 330), "Module name: Research Methods in Social Sciences", fontsize=10)
    page.insert_text((50, 350), "Module lead's name: John Connor", fontsize=10)
    page.insert_text((50, 370), "Module Lead: Sarah Connor", fontsize=10)
    page.insert_text((50, 390), "Submission location: Turnitin VLE online portal", fontsize=10)
    page.insert_text((50, 410), "Submit to: Blackboard submission folder", fontsize=10)
    page.insert_text((50, 425), "Academic Year 2025-26", fontsize=10)
    page.insert_text((50, 435), "Module Code: BM414", fontsize=10)

    # 4. Middle of the page office address block (should be globally detected/redacted)
    office_text = "HEAD OFFICE\n7 Acorn Business Park\nCommercial Gate, Nottingham\nNG18 1EX"
    page.insert_textbox(fitz.Rect(50, 600, 400, 680), office_text, fontsize=10)


    # Insert university_2 logo
    if os.path.exists(logo_path):
        page.insert_image(fitz.Rect(50, 430, 200, 580), filename=logo_path)
        
    doc.save(filename)
    doc.close()
    print(f"Created test PDF: {filename}")

def create_test_pptx(filename: str):
    prs = pptx.Presentation()
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title and Content
    slide.shapes.title.text = "Oxford Orientation Slide"
    
    txBox = slide.shapes.add_textbox(pptx.util.Inches(1), pptx.util.Inches(2), pptx.util.Inches(8), pptx.util.Inches(2))
    tf = txBox.text_frame
    
    p = tf.add_paragraph()
    p.text = "Roll number of student: ROLL12345"
    p = tf.add_paragraph()
    p.text = "Enquiries: admin@ox.ac.uk or +44-1865-270000"
    p = tf.add_paragraph()
    p.text = "Submitted on Oct 12, 2025 at 10:00 AM"
    
    p = tf.add_paragraph()
    p.text = "Name/Signed: Nargisa Simansone"
    p = tf.add_paragraph()
    p.text = "Module leader: Adeeba Ahmad"
    p = tf.add_paragraph()
    p.text = "Module lead's name: John Connor"
    p = tf.add_paragraph()
    p.text = "Module Lead: Sarah Connor"
    p = tf.add_paragraph()
    p.text = "Academic Year 2025-26"
    p = tf.add_paragraph()
    p.text = "Module Code: BM414"

    p = tf.add_paragraph()
    p.text = "Submission location: Turnitin VLE online portal"
    p = tf.add_paragraph()
    p.text = "Submit to: Blackboard submission folder"
    
    # Add university_3 logo
    logo_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "logos", "university_3.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, pptx.util.Inches(1), pptx.util.Inches(4.5), width=pptx.util.Inches(2.5))
        
    prs.save(filename)
    print(f"Created test PPTX: {filename}")
 
def run_tests():
    temp_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "temp_test")
    os.makedirs(temp_dir, exist_ok=True)
    
    in_docx = os.path.join(temp_dir, "input.docx")
    out_docx = os.path.join(temp_dir, "output.docx")
    in_pdf = os.path.join(temp_dir, "input.pdf")
    out_pdf = os.path.join(temp_dir, "output.pdf")
    in_pptx = os.path.join(temp_dir, "input.pptx")
    out_pptx = os.path.join(temp_dir, "output.pptx")
    
    create_test_docx(in_docx)
    create_test_pdf(in_pdf)
    create_test_pptx(in_pptx)
    
    print("\n--- Running Redaction Processes ---")
    
    # 1. Test DOCX
    print("Testing DOCX processing...")
    process_docx(in_docx, out_docx)
    doc = docx.Document(out_docx)
    docx_text = ""
    for p in doc.paragraphs:
        docx_text += p.text + "\n"
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                docx_text += cell.text + "\n"
    # Check headers/footers
    for section in doc.sections:
        if section.header:
            for p in section.header.paragraphs:
                docx_text += p.text + "\n"
        if section.footer:
            for p in section.footer.paragraphs:
                docx_text += p.text + "\n"
                
    assert "ST12345" not in docx_text, "DOCX: ST12345 not redacted"
    assert "REG2024001" not in docx_text, "DOCX: REG2024001 not redacted"
    assert "harvard.edu" in docx_text, "DOCX: harvard.edu domain should be preserved"
    assert "Harvard" in docx_text, "DOCX: Harvard reference should be preserved"
    assert "BNU" not in docx_text, "DOCX: BNU issuing university alias not redacted"
    assert "bucks.ac.uk" not in docx_text, "DOCX: bucks.ac.uk issuing university domain not redacted"
    assert "Buckinghamshire" not in docx_text, "DOCX: Buckinghamshire issuing university name not redacted"
    assert "+1 (123) 456-7890" not in docx_text, "DOCX: phone not redacted"
    assert "SW1A 1AA" not in docx_text, "DOCX: postcode not redacted"
    assert "19th" not in docx_text, "DOCX: submission date (19th June) not redacted"
    assert "June 2026" not in docx_text, "DOCX: submission date (June 2026) not redacted"
    assert "10th" not in docx_text, "DOCX: target feedback date (10th July) not redacted"
    assert "July 2026" not in docx_text, "DOCX: target feedback date (July 2026) not redacted"
    assert "21st" not in docx_text, "DOCX: resit date (21st August) not redacted"
    assert "August 2026" not in docx_text, "DOCX: resit date (August 2026) not redacted"
    assert "4:00pm" not in docx_text, "DOCX: military time 4:00pm not redacted"
    assert "1600hrs" not in docx_text, "DOCX: military time 1600hrs not redacted"
    assert "12 June" not in docx_text, "DOCX: military time date (12 June) not redacted"
    assert "25th October 2026" not in docx_text, "DOCX: Draft Submission date not redacted"
    assert "26th October 2026" not in docx_text, "DOCX: Draft Submission (Mandatory) date not redacted"
    assert "27th October 2026" not in docx_text, "DOCX: Formative Submission date not redacted"
    assert "28th October 2026" not in docx_text, "DOCX: Draft Submission (Mandatory) paragraph date not redacted"
    assert "29th October 2026" not in docx_text, "DOCX: Date and Time of Submission date not redacted"
    assert "30th October 2026" not in docx_text, "DOCX: Submission Date & Time date not redacted"
    assert "31st October 2026" not in docx_text, "DOCX: Date and Time of Submission paragraph date not redacted"
    assert "Adeeba" not in docx_text, "DOCX: Adeeba not redacted"
    assert "Ahmad" not in docx_text, "DOCX: Ahmad not redacted"
    assert "Nargisa" not in docx_text, "DOCX: Nargisa not redacted"
    assert "Simansone" not in docx_text, "DOCX: Simansone not redacted"
    assert "John Smith" not in docx_text, "DOCX: John Smith not redacted"
    assert "John Connor" not in docx_text, "DOCX: John Connor not redacted"
    assert "Sarah Connor" not in docx_text, "DOCX: Sarah Connor not redacted"
    assert "Turnitin VLE online portal" not in docx_text, "DOCX: Submission location value not redacted"
    assert "Blackboard submission folder" not in docx_text, "DOCX: Submit-to value not redacted"
    
    # Assertions for Layout Types 1-4
    assert "Turnitin Submission Portal" not in docx_text, "DOCX: Turnitin Submission Portal not redacted"
    assert "Option A: Presentation" not in docx_text, "DOCX: Option A: Presentation not redacted"
    assert "Option B: Turnitin Submission" not in docx_text, "DOCX: Option B: Turnitin Submission not redacted"
    assert "Option A + Option B + Option C" not in docx_text, "DOCX: Option A + Option B + Option C not redacted"

    # Assertions for Free-text multi-paragraph options
    assert "Option A: Online Upload" not in docx_text, "DOCX: Option A: Online Upload not redacted"
    assert "Option B: In-person Hand-in" not in docx_text, "DOCX: Option B: In-person Hand-in not redacted"
    assert "Prof. Green" not in docx_text, "DOCX: Prof. Green not redacted"
    assert "OCR Failure Test" not in docx_text, "DOCX: OCR Failure Test (fuzzy match) not redacted"

    # Assertions for Academic exclusions being preserved (NOT redacted)
    assert "Managing Innovation" in docx_text, "DOCX: 'Managing Innovation' (academic exclusion) was incorrectly redacted"
    assert "Reference List" in docx_text, "DOCX: 'Reference List' (academic exclusion) was incorrectly redacted"
    assert "Recommended Reading" in docx_text, "DOCX: 'Recommended Reading' (academic exclusion) was incorrectly redacted"
    assert "Learning Outcomes" in docx_text, "DOCX: 'Learning Outcomes' (academic exclusion) was incorrectly redacted"
    assert "Module Guide" in docx_text, "DOCX: 'Module Guide' (academic exclusion) was incorrectly redacted"
    assert "Assessment Brief" in docx_text, "DOCX: 'Assessment Brief' (academic exclusion) was incorrectly redacted"
    assert "Confidentiality" in docx_text, "DOCX: 'Confidentiality' section was incorrectly redacted"
    assert "Academic Integrity" in docx_text, "DOCX: 'Academic Integrity' section was incorrectly redacted"

    # Assertions for human names being redacted
    assert "Claire Ngo" not in docx_text, "DOCX: Claire Ngo not redacted"
    assert "Sarah Johnson" not in docx_text, "DOCX: Sarah Johnson not redacted"
    assert "Michael Brown" not in docx_text, "DOCX: Michael Brown not redacted"
    assert "Academic Year" not in docx_text, "DOCX: Academic Year not redacted"
    assert "BM414" not in docx_text, "DOCX: BM414 not redacted"


    print("DOCX Redaction Verification: SUCCESS")

    
    # 2. Test PDF
    print("Testing PDF processing...")
    process_pdf(in_pdf, out_pdf)
    pdf_doc = fitz.open(out_pdf)
    pdf_text = ""
    for page in pdf_doc:
        pdf_text += page.get_text()
    pdf_doc.close()
    
    assert "APP99881" not in pdf_text, "PDF: APP99881 not redacted"
    assert "registry@stanford.edu" not in pdf_text, "PDF: Email not redacted"
    assert "stanford.edu" not in pdf_text, "PDF: stanford.edu not redacted"
    assert "Stanford" not in pdf_text, "PDF: Stanford reference not redacted"
    assert "650-123-4567" not in pdf_text, "PDF: Phone not redacted"
    assert "94305" not in pdf_text, "PDF: ZIP code not redacted"
    assert "14:30" not in pdf_text, "PDF: target feedback time not redacted"
    assert "21st" not in pdf_text, "PDF: resit date (21st August) not redacted"
    assert "August 2026" not in pdf_text, "PDF: resit date (August 2026) not redacted"
    assert "4:00pm" not in pdf_text, "PDF: military time 4:00pm not redacted"
    assert "1600hrs" not in pdf_text, "PDF: military time 1600hrs not redacted"
    assert "12 June" not in pdf_text, "PDF: military time date (12 June) not redacted"
    assert "Adeeba" not in pdf_text, "PDF: Adeeba not redacted"
    assert "Ahmad" not in pdf_text, "PDF: Ahmad not redacted"
    assert "Nargisa" not in pdf_text, "PDF: Nargisa not redacted"
    assert "Simansone" not in pdf_text, "PDF: Simansone not redacted"
    assert "John Connor" not in pdf_text, "PDF: John Connor not redacted"
    assert "Sarah Connor" not in pdf_text, "PDF: Sarah Connor not redacted"
    assert "Turnitin VLE online portal" not in pdf_text, "PDF: Submission location value not redacted"
    assert "Blackboard submission folder" not in pdf_text, "PDF: Submit-to value not redacted"

    # Assertions for Programme and Module name not being redacted in PDF
    assert "Programme" in pdf_text, "PDF: Programme label got redacted"
    assert "Module name" in pdf_text, "PDF: Module name label got redacted"
    assert "Health and Social Care" in pdf_text, "PDF: Programme value got redacted"
    assert "Research Methods" in pdf_text, "PDF: Module name value got redacted"
    assert "Academic Year" not in pdf_text, "PDF: Academic Year not redacted"
    assert "BM414" not in pdf_text, "PDF: BM414 not redacted"

    # Assertions for header branding & logo removal metrics
    from python.redaction.redaction_audit import RedactionAudit
    pdf_summary = RedactionAudit.generate_summary("input.pdf")
    assert pdf_summary.get("header_images_removed", 0) > 0, f"Expected header_images_removed > 0, got {pdf_summary.get('header_images_removed')}"
    assert pdf_summary.get("address_blocks_removed", 0) > 0, f"Expected address_blocks_removed > 0, got {pdf_summary.get('address_blocks_removed')}"
    assert pdf_summary.get("contact_blocks_removed", 0) > 0, f"Expected contact_blocks_removed > 0, got {pdf_summary.get('contact_blocks_removed')}"

    # Verify global address block is redacted
    assert "Acorn Business Park" not in pdf_text, "PDF: Global address block was not redacted"
    assert "NG18 1EX" not in pdf_text, "PDF: Global address postcode was not redacted"


    print("PDF Redaction Verification: SUCCESS")
    
    # 3. Test PPTX
    print("Testing PPTX processing...")
    process_pptx(in_pptx, out_pptx)
    prs = pptx.Presentation(out_pptx)
    pptx_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                pptx_text += shape.text_frame.text + "\n"
                
    assert "ROLL12345" not in pptx_text, "PPTX: ROLL12345 not redacted"
    assert "admin@ox.ac.uk" not in pptx_text, "PPTX: Email not redacted"
    assert "ox.ac.uk" not in pptx_text, "PPTX: ox.ac.uk not redacted"
    assert "Oxford" not in pptx_text, "PPTX: Oxford reference not redacted"
    assert "+44-1865-270000" not in pptx_text, "PPTX: Phone not redacted"
    assert "Oct 12" not in pptx_text, "PPTX: submission date not redacted"
    assert "Adeeba" not in pptx_text, "PPTX: Adeeba not redacted"
    assert "Ahmad" not in pptx_text, "PPTX: Ahmad not redacted"
    assert "Nargisa" not in pptx_text, "PPTX: Nargisa not redacted"
    assert "Simansone" not in pptx_text, "PPTX: Simansone not redacted"
    assert "John Connor" not in pptx_text, "PPTX: John Connor not redacted"
    assert "Sarah Connor" not in pptx_text, "PPTX: Sarah Connor not redacted"
    assert "Turnitin VLE online portal" not in pptx_text, "PPTX: Submission location value not redacted"
    assert "Blackboard submission folder" not in pptx_text, "PPTX: Submit-to value not redacted"
    assert "Academic Year" not in pptx_text, "PPTX: Academic Year not redacted"
    assert "BM414" not in pptx_text, "PPTX: BM414 not redacted"

    print("PPTX Redaction Verification: SUCCESS")
    
    print("\n--- ALL TESTS COMPLETED SUCCESSFULLY ---")

if __name__ == "__main__":
    run_tests()
